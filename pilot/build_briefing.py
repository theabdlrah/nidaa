#!/usr/bin/env python3
"""Build 'Nidaa Founder Briefing v1' (.pptx) — synthesis of existing project artifacts.
No new research; every figure/claim is grounded in files already written in pilot/.
Run: python build_briefing.py  ->  produces Nidaa-Founder-Briefing-v1.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- palette ----
NAVY   = RGBColor(0x1C, 0x2B, 0x39)   # deep background
TEAL   = RGBColor(0x15, 0x9A, 0x9A)   # accent (offline-first / structured)
SAND   = RGBColor(0xF2, 0xE9, 0xD8)   # light bg / text on dark
INK     = RGBColor(0x22, 0x2A, 0x33)   # dark text on light
GREY    = RGBColor(0x6B, 0x72, 0x80)
RED     = RGBColor(0xC0, 0x39, 0x2B)
GREEN   = RGBColor(0x2E, 0x8B, 0x57)
AMBER   = RGBColor(0xD9, 0x8A, 0x10)
BLUE    = RGBColor(0x2D, 0x6C, 0xDF)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

SW, SH = Inches(13.333), Inches(7.5)   # 16:9

prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]

def slide():
    return prs.slides.add_slide(BLANK)

def rect(s, x, y, w, h, fill, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp

def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=4, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is list of (text, size, bold, color, italic)"""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        for (t, sz, b, col, *rest) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.bold = b
            r.font.color.rgb = col
            r.font.name = "Calibri"
            if rest and rest[0]:
                r.font.italic = True
    return tb

def header(s, kicker, title, num):
    rect(s, 0, 0, SW, Inches(1.15), NAVY)
    rect(s, 0, Inches(1.15), SW, Pt(3), TEAL)
    txt(s, Inches(0.55), Inches(0.16), Inches(11), Inches(0.32),
        [[(kicker.upper(), 12, True, TEAL)]])
    txt(s, Inches(0.55), Inches(0.46), Inches(11.2), Inches(0.62),
        [[(title, 24, True, WHITE)]])
    txt(s, Inches(12.2), Inches(0.42), Inches(0.9), Inches(0.5),
        [[(num, 13, True, GREY)]], align=PP_ALIGN.RIGHT)

def footer(s):
    txt(s, Inches(0.55), Inches(7.08), Inches(9), Inches(0.32),
        [[("Nidaa (نداء) — Founder Briefing v1 · 2026-07-17 · Syria/Gaza offline-first aid board", 9, False, GREY)]])

def bullet(s, x, y, w, h, items, size=14, gap=7, color=INK, lead=TEAL):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.line_spacing = 1.02
        if isinstance(it, tuple):
            head, body = it
            r = p.add_run(); r.text = "▪ "; r.font.size = Pt(size); r.font.color.rgb = lead; r.font.bold = True
            r2 = p.add_run(); r2.text = head; r2.font.size = Pt(size); r2.font.bold = True; r2.font.color.rgb = color
            r3 = p.add_run(); r3.text = "  " + body; r3.font.size = Pt(size); r3.font.color.rgb = color
        else:
            r = p.add_run(); r.text = "▪ "; r.font.size = Pt(size); r.font.color.rgb = lead; r.font.bold = True
            r2 = p.add_run(); r2.text = it; r2.font.size = Pt(size); r2.font.color.rgb = color
        for rr in p.runs:
            rr.font.name = "Calibri"
    return tb

def chip(s, x, y, w, label, fill, tcol=WHITE, h=Inches(0.36), size=11):
    c = rect(s, x, y, w, h, fill)
    tf = c.text_frame; tf.word_wrap = True
    tf.margin_top=Pt(1); tf.margin_bottom=Pt(1)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label; r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = tcol
    r.font.name = "Calibri"
    return c

# ============================================================ SLIDE 1
s = slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(3.05), SW, Pt(3), TEAL)
txt(s, Inches(0.9), Inches(1.55), Inches(11.5), Inches(0.5),
    [[("نداء  ·  NIDAA", 30, True, TEAL)]])
txt(s, Inches(0.9), Inches(2.15), Inches(11.5), Inches(1.0),
    [[("Offline-first, Arabic-RTL community needs & services board", 26, True, WHITE)]])
txt(s, Inches(0.9), Inches(3.25), Inches(11.5), Inches(1.6),
    [[("A humanitarian aid-coordination board built for how Syria & Gaza actually are in 2025–26:", 16, False, SAND)],
     [("~64% offline · intermittent & expensive connectivity · damaged infrastructure · information dies exactly when it is needed most.", 16, False, SAND, True)]],
    space_after=8)
txt(s, Inches(0.9), Inches(5.2), Inches(11.5), Inches(0.4),
    [[("Founder Briefing v1  ·  Abdul Rahman  ·  2026-07-17", 14, True, TEAL)]])
txt(s, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.5),
    [[("Status: prototype built & verified · evidence campaign in progress · not yet a deployed service", 12, False, GREY)]])

# ============================================================ SLIDE 2 — what is
s = slide(); header(s, "The product", "1 · What is Nidaa?", "01 / 15")
bullet(s, Inches(0.55), Inches(1.5), Inches(7.0), Inches(4.6), [
    ("Offline-first by design.", "The device is the source of truth — new posts save locally (IndexedDB) and sync only when a connection appears."),
    ("Works with zero connectivity.", "Service-worker app shell means the board still loads and functions after the first visit."),
    ("Arabic-default, full RTL.", "With an English toggle — built for the people on the ground, not for HQ."),
    ("Structured, verifiable entries.", "Needs & offers (medical, food, water, shelter, transport…) by city; NGOs mark entries verified."),
    ("Gaza-first, Syria-second.", "One codebase, region-tagged (gza | syr). Gaza is the primary deployment target."),
], size=15, gap=10)
# right panel = stack
rect(s, Inches(7.9), Inches(1.5), Inches(4.85), Inches(4.85), SAND)
txt(s, Inches(8.1), Inches(1.62), Inches(4.4), Inches(0.4),
    [[("What it actually does", 15, True, NAVY)]])
bullet(s, Inches(8.1), Inches(2.05), Inches(4.45), Inches(4.2), [
    "Browse community needs & offers by city",
    "Post a new entry even when offline (queues, pushes on reconnect)",
    "NGOs / admins mark entries verified (auth-gated in prod)",
    "Live online/offline + pending-sync indicators",
    "Pluggable real seed data (HDX health/education facilities)",
], size=13, gap=9, color=INK)
footer(s)

# ============================================================ SLIDE 3 — why problem matters
s = slide(); header(s, "The problem", "2 · Why does this matter?", "02 / 15")
bullet(s, Inches(0.55), Inches(1.5), Inches(12.2), Inches(2.4), [
    ("Connectivity is the failure mode.", "In Syria ~64% are offline and internet is slow/intermittent/expensive; in Gaza, internet AND electricity scarcity are the binding constraints."),
    ("Incumbent tools fail exactly when needed.", "WhatsApp / Telegram / Facebook work only with a connection — they break during the outages when aid info matters most."),
    ("Information fragmentation has real cost.", "A resident interview (Adam, 2026-07-14): a connectivity outage cut his family's access to food-distribution information and to digital payments."),
], size=15, gap=11)
# three stat chips
chip(s, Inches(0.55), Inches(4.35), Inches(3.9), "≈ 64%  offline (Syria)", NAVY, size=14)
chip(s, Inches(4.7), Inches(4.35), Inches(3.9), "~36%  internet penetration", TEAL, size=14)
chip(s, Inches(8.85), Inches(4.35), Inches(3.9), "0 tools fit offline+Gaza", AMBER, size=14)
txt(s, Inches(0.55), Inches(5.05), Inches(12.2), Inches(1.6),
    [[("The design principle that follows:", 14, True, NAVY)],
     [("\u201cthe network will fail, and the tool must be correct anyway.\u201d  Every layer (storage, UI, verification, sync) preserves offline-first behavior — Nidaa degrades gracefully instead of going dark.", 15, False, INK)]],
    space_after=8)
footer(s)

# ============================================================ SLIDE 4 — timeline
s = slide(); header(s, "How we got here", "3 · Project evolution timeline", "03 / 15")
rows = [
    ("Jul 14", "Resident interview (Adam)", "A2/A3/A5 strengthened: info + offline matter; incumbents fail on outages.", GREEN),
    ("Jul 15", "Red-team + pivot to coordinator-led", "Hypothesis shifts: primary user = coordinator / info officer, not lone beneficiary.", BLUE),
    ("Jul 15", "A6 splits into A6a / A6b", "Formal orgs structure info (early support); grassroots coordinators = the open, existential question.", AMBER),
    ("Jul 15", "Public-evidence matrix (Tier-1)", "9 sources retrieved + URL-verified; convergence becomes the strongest signal.", TEAL),
    ("Jul 15–16", "Outreach wave sent", "17 targets contacted / engaged incl. Abdelrahman (Wa7at), OCHA, WCK, YAP, Mones (MSF).", GREY),
    ("Jul 17", "A6b Resolution Campaign", "RESOLVED-AS-CONDITIONAL: structure survives only with owner/mandate + near-zero burden.", NAVY),
]
y = Inches(1.45)
for d, t, desc, col in rows:
    rect(s, Inches(0.55), y, Inches(1.35), Inches(0.78), col)
    txt(s, Inches(0.55), y+Inches(0.13), Inches(1.35), Inches(0.5),
        [[(d, 12, True, WHITE)]], align=PP_ALIGN.CENTER)
    rect(s, Inches(1.95), y, Inches(10.8), Inches(0.78), SAND)
    txt(s, Inches(2.1), y+Inches(0.06), Inches(10.5), Inches(0.34),
        [[(t, 13.5, True, NAVY)]])
    txt(s, Inches(2.1), y+Inches(0.40), Inches(10.5), Inches(0.34),
        [[(desc, 11.5, False, INK)]])
    y += Inches(0.92)
footer(s)

# ============================================================ SLIDE 5 — current hypothesis
s = slide(); header(s, "The thesis today", "4 · Current hypothesis", "04 / 15")
rect(s, Inches(0.55), Inches(1.5), Inches(12.2), Inches(1.5), NAVY)
txt(s, Inches(0.8), Inches(1.66), Inches(11.7), Inches(1.2),
    [[("Nidaa may be a coordinator-led information-management & coordination tool, not a civilian-first platform.", 18, True, WHITE)],
     [("The user is the coordinator / info officer who already \u201ckeeps information\u201d — Nidaa makes that job near-zero-burden and embeds inside their existing channel.", 14, False, SAND)]],
    space_after=10)
txt(s, Inches(0.55), Inches(3.2), Inches(12.2), Inches(0.4),
    [[("Why this hypothesis (and why it is still falsifiable):", 14, True, NAVY)]])
bullet(s, Inches(0.55), Inches(3.65), Inches(12.2), Inches(2.8), [
    ("Tool research showed mature systems always have an owner.", "Ushahidi→deployment owner, HOTOSM→validators, ReliefWeb→OCHA editors, HDX→org data managers."),
    ("A7 split accordingly.", "A7a = coordinators publish (now less risky); A7b = beneficiaries use (still untested)."),
    ("A6b then sharpened the shape.", "Structure survives only when (a) there is an owner/mandate and (b) burden is near-zero and lives in the existing workflow."),
    ("Reality still gets the vote.", "This is a working hypothesis, not a conclusion — every pending coordinator conversation is a falsification test."),
], size=14, gap=9)
footer(s)

# ============================================================ SLIDE 6 — assumption map
s = slide(); header(s, "The assumption map", "5 · What we are testing (A1–A7, A6a, A6b, A6-trust)", "05 / 15")
cols = [
    ("A1", "Coordination is a priority problem", "Paused — bounded", AMBER),
    ("A2", "Info/matching is a bottleneck", "Strengthened", GREEN),
    ("A3", "Offline capability materially matters", "Strengthened", GREEN),
    ("A4", "Verification materially matters", "Untested", GREY),
    ("A5", "Existing tools insufficient", "Strengthened", GREEN),
    ("A6-trust", "Orgs trust designated verifiers", "Untested", GREY),
    ("A7", "Communities adopt a new workflow", "Untested", GREY),
    ("A6a", "Formal orgs maintain structured info", "Early support", TEAL),
    ("A6b", "Grassroots coordinators maintain structured info", "Conditional", BLUE),
]
x0 = Inches(0.55); y0 = Inches(1.5)
cw = Inches(4.0); ch = Inches(1.18); gx = Inches(0.13); gy = Inches(0.13)
for i, (tag, label, st, col) in enumerate(cols):
    cx = x0 + (cw + gx) * (i % 3)
    cy = y0 + (ch + gy) * (i // 3)
    rect(s, cx, cy, cw, ch, SAND)
    rect(s, cx, cy, Inches(1.25), ch, col)
    txt(s, cx, cy+Inches(0.30), Inches(1.25), Inches(0.6),
        [[(tag, 19, True, WHITE)]], align=PP_ALIGN.CENTER)
    txt(s, cx+Inches(1.38), cy+Inches(0.10), cw-Inches(1.5), Inches(0.62),
        [[(label, 12.5, True, NAVY)]])
    txt(s, cx+Inches(1.38), cy+Inches(0.74), cw-Inches(1.5), Inches(0.36),
        [[(st, 12, True, col)]])
footer(s)

# ============================================================ SLIDE 7 — evidence so far
s = slide(); header(s, "Evidence collected", "6 · What the evidence shows so far", "06 / 15")
bullet(s, Inches(0.55), Inches(1.5), Inches(12.2), Inches(3.0), [
    ("Tier-1 public evidence — 9 sources retrieved & URL-verified.", "Incl. ACAPS (ERRs Sudan), White Helmets reports, GTS, CDAC, HOTOSM, MapAction, ReliefWeb, SARC, WFP — converged where independent sources agree."),
    ("Tier-2 field — 1 resident interview (Adam, Jul 14).", "Strengthened A2/A3/A5 with first-hand outage failure: lost access to food-distribution info & digital payments."),
    ("1 partial reply (Abdelrahman Saleh, Wa7at).", "Named internet + electricity scarcity as major challenges; detailed questionnaire pending."),
    ("A6b campaign — 5 evidence items + 4 contradictions logged.", "Population-split now evidence-backed, not guessed."),
], size=14, gap=10)
# convergence strip
txt(s, Inches(0.55), Inches(4.75), Inches(12), Inches(0.35),
    [[("Convergence status (independent-source agreement):", 13, True, NAVY)]])
strip = [("A3", "Emerging", GREEN), ("A2/A5", "Weak (1 interview)", AMBER),
         ("A6a", "Early support", TEAL), ("A6b", "Conditional", BLUE),
         ("A1", "Bounded", AMBER), ("A4/A6-trust/A7", "Untested", GREY)]
sx = Inches(0.55)
for tag, st, col in strip:
    w = Inches(2.35)
    chip(s, sx, Inches(5.15), w, tag + ": " + st, col, size=10.5)
    sx += w + Inches(0.04)
txt(s, Inches(0.55), Inches(5.75), Inches(12.2), Inches(1.0),
    [[("No assumption has been falsified. Pivot Rule NOT triggered. A3 is the only Emerging-Convergence signal so far.", 13, False, INK, True)]])
footer(s)

# ============================================================ SLIDE 8 — major findings
s = slide(); header(s, "Findings", "7 · Major findings", "07 / 15")
bullet(s, Inches(0.55), Inches(1.5), Inches(12.2), Inches(4.8), [
    ("The naive standalone-board model is dead.", "No informal-SUCCESS case exists for a loose coordinator spontaneously maintaining a register."),
    ("But Nidaa is not dead either.", "White Helmets prove volunteers DO structure under extreme pressure — when reporting is mandated, trained, and accountable."),
    ("The real pattern is a moderator, not a yes/no.", "Structure survives where an OWNER + MANDATE exist and burden is near-zero; it collapses without them."),
    ("Informal coordination defaults to chat, not structure.", "ERRs Sudan: the \u201croom\u201d is the group chat; documentation is named as a GAP to be built, not a practice."),
    ("Structuring is done BY/THROUGH a facilitating function.", "GTS, CDAC, White Helmets all imply a function owns the information — rarely the crowd itself."),
], size=14.5, gap=10)
footer(s)

# ============================================================ SLIDE 9 — Reality A vs B
s = slide(); header(s, "The fork", "8 · Reality A vs Reality B", "08 / 15")
# A panel
rect(s, Inches(0.55), Inches(1.5), Inches(6.0), Inches(4.7), SAND)
rect(s, Inches(0.55), Inches(1.5), Inches(6.0), Inches(0.6), GREEN)
txt(s, Inches(0.55), Inches(1.58), Inches(6.0), Inches(0.5),
    [[("REALITY A — structure survives", 16, True, WHITE)]], align=PP_ALIGN.CENTER)
bullet(s, Inches(0.75), Inches(2.2), Inches(5.6), Inches(3.9), [
    "Spreadsheet / register / list maintained",
    "A named focal point or info officer",
    "Status tracking (who-has-what, what's done)",
    "Structured, regular, written updates",
    "Evidence: White Helmets (field + monthly reports under bombardment), GTS, CDAC, A6a orgs",
], size=13, gap=10, color=INK)
# B panel
rect(s, Inches(6.75), Inches(1.5), Inches(6.0), Inches(4.7), SAND)
rect(s, Inches(6.75), Inches(1.5), Inches(6.0), Inches(0.6), AMBER)
txt(s, Inches(6.75), Inches(1.58), Inches(6.0), Inches(0.5),
    [[("REALITY B — communication survives, structure doesn't", 15, True, WHITE)]], align=PP_ALIGN.CENTER)
bullet(s, Inches(6.95), Inches(2.2), Inches(5.6), Inches(3.9), [
    "Calls only / WhatsApp only",
    "Verbal coordination",
    "No maintained records",
    "Records started but abandoned (\u201cthe doc died after 3 days\u201d)",
    "Evidence: ERRs Sudan (chat-native, docs-as-gap)",
], size=13, gap=10, color=INK)
txt(s, Inches(0.55), Inches(6.35), Inches(12.2), Inches(0.7),
    [[("Reality B = Discovery, not failure. It redirects Nidaa from \u201ca separate board\u201d toward \u201cembed inside the existing workflow / be the facilitating function\u201d — a better-shaped product.", 13.5, True, NAVY)]])
footer(s)

# ============================================================ SLIDE 10 — why A6b changed the project
s = slide(); header(s, "The pivot", "9 · Why A6b changed the project", "09 / 15")
txt(s, Inches(0.55), Inches(1.45), Inches(12.2), Inches(0.8),
    [[("Before A6b: a guess that informal coordinators \u201cprobably won't structure\u201d (Unproven, leaning Weak).", 15, False, INK)],
     [("After A6b: a precise, testable CONDITIONAL — and the \u201cwhy\u201d is now grounded in real fetched sources, not speculation.", 15, True, NAVY)]],
    space_after=8)
bullet(s, Inches(0.55), Inches(2.7), Inches(12.2), Inches(3.4), [
    ("It killed the naive version of A6b.", "No evidence that a standalone coordinator-led board gets maintained by a loose mutual-aid chat."),
    ("It replaced it with a sharp conditional.", "\u201cStructure holds IF owner/mandate + near-zero burden + embedded in existing channel.\u201d"),
    ("It reshaped the product direction.", "Toward \u201cembed in the workflow / be the facilitating function\u201d — i.e. Reality B as a design pivot, now evidence-backed."),
    ("It is a genuine falsification result.", "This is exactly what the campaign was for: it invalidated one claim and produced a better, narrower one."),
], size=14.5, gap=10)
chip(s, Inches(0.55), Inches(6.25), Inches(12.2), "A6b is now RESOLVED-AS-CONDITIONAL (Medium confidence) — the defining result of the project so far.", NAVY, size=14, h=Inches(0.55))
footer(s)

# ============================================================ SLIDE 11 — highest-risk remaining
s = slide(); header(s, "Risk", "10 · Highest-risk assumptions remaining", "10 / 15")
rect(s, Inches(0.55), Inches(1.5), Inches(12.2), Inches(1.25), RED)
txt(s, Inches(0.8), Inches(1.62), Inches(11.7), Inches(1.0),
    [[("A7 — Communities adopt a new workflow  (UNTESTED, Low confidence)", 17, True, WHITE)],
     [("Most existential AND least evidenced. The A6b result raises the bar: structure survives only if embedded with near-zero burden — that is an adoption/UX problem, not a capability one.", 13, False, SAND)]],
    space_after=6)
bullet(s, Inches(0.55), Inches(2.95), Inches(12.2), Inches(3.4), [
    ("A4 — Verification materially matters.", "Untested (3 conceptual public sources, 0 first-hand verification pain)."),
    ("A6-trust — Orgs trust designated verifiers.", "Untested. Weak early counter-signal: beneficiaries expect info pushed (WhatsApp), not pulled from a board."),
    ("A1 — Coordination is a priority problem.", "Paused at a bounded conclusion (binding where capacity+access exist; secondary where access/security dominate) — favourably bounded, but not closed."),
    ("A6b sub-claims still Low.", "No first-hand informal-SUCCESS or informal-COLLAPSE event yet — needs a real coordinator story."),
], size=14, gap=10)
footer(s)

# ============================================================ SLIDE 12 — outreach pipeline
s = slide(); header(s, "Pipeline", "11 · Current outreach pipeline", "11 / 15")
txt(s, Inches(0.55), Inches(1.4), Inches(12.2), Inches(0.4),
    [[("17 targets approached/engaged · Row-1 (first real coordinator story) is the live gate", 14, True, NAVY)]])
# two columns: pending / warm
rect(s, Inches(0.55), Inches(1.95), Inches(6.0), Inches(4.6), SAND)
txt(s, Inches(0.75), Inches(2.05), Inches(5.6), Inches(0.4),
    [[("Warm / highest-value", 14, True, NAVY)]])
bullet(s, Inches(0.75), Inches(2.5), Inches(5.6), Inches(3.9), [
    ("Mones (MSF ops)", "Contacted — #1 Row-1 candidate, pending"),
    ("Abdelrahman Saleh (Wa7at)", "Replied — detailed questionnaire pending (Gaza)"),
    ("Fatma (health-team referral)", "Contacted — potential referral into Gaza ground coords"),
    ("Omar / Shaima / local leads", "Queued — informal-coordinator signal"),
], size=13, gap=11, color=INK)
rect(s, Inches(6.75), Inches(1.95), Inches(6.0), Inches(4.6), SAND)
txt(s, Inches(6.95), Inches(2.05), Inches(5.6), Inches(0.4),
    [[("Org / institutional", 14, True, NAVY)]])
bullet(s, Inches(6.95), Inches(2.5), Inches(5.6), Inches(3.9), [
    ("OCHA (Husam, Alaa, Georgios)", "Contacted — coordination-insight targets"),
    ("WCK (Ibrahim), YAP Palestine", "Contacted — 07-16"),
    ("iMMAP / REACH / HOTOSM", "Verification + structured-info angles"),
    ("White Helmets / ERRs / Sameer", "SUCCESS & FAILURE-case targets (next wave)"),
], size=13, gap=11, color=INK)
txt(s, Inches(0.55), Inches(6.65), Inches(12.2), Inches(0.4),
    [[("Discipline: outreach status is NOT evidence. A reply alone never counts — only a real story routed via the evidence pipeline.", 12, False, GREY, True)]])
footer(s)

# ============================================================ SLIDE 13 — what would kill
s = slide(); header(s, "Stress test", "12 · What would kill Nidaa?", "12 / 15")
bullet(s, Inches(0.55), Inches(1.5), Inches(12.2), Inches(4.9), [
    ("A7 falsified at scale.", "Beneficiaries never open such boards / rely solely on WhatsApp forwards / distrust unknown sources — and coordinators won't adopt the workflow either."),
    ("A1 strongly falsified by multiple conversations.", "Orgs rank coordination below fuel/access/security; \u201ccoordination is fine, X is the real problem.\u201d → triggers the Pivot Rule."),
    ("A6-trust collapses.", "Orgs refuse to post into a shared board citing liability / false-info harm / \u201cwe control our own comms.\u201d"),
    ("The informal-COLLAPSE is total and unmediated.", "Real coordinators try a register and it dies under surge/connectivity loss with no facilitating function available — and no embeddable channel exists."),
    ("Adoption burden is not actually near-zero.", "If embedding in the workflow still requires meaningful admin effort, the A6b conditional fails in practice."),
], size=14.5, gap=11)
footer(s)

# ============================================================ SLIDE 14 — what would validate
s = slide(); header(s, "Stress test", "13 · What would validate Nidaa?", "13 / 15")
bullet(s, Inches(0.55), Inches(1.5), Inches(12.2), Inches(4.9), [
    ("Row-1 lands — a real coordinator story answering all six A6 questions.", "Either a SUCCESS (structured maintenance held) or a COLLAPSE (dropped under load). Either is evidence; the milestone is the story, not a report."),
    ("A7 gets a first real signal.", "A coordinator or community adopts the workflow because it solved a concrete problem — not because they liked the idea."),
    ("A6-trust earns one named verifier.", "An org with hierarchy names a person/body they'd accept as verifier, or demonstrably welcomes a verifier role."),
    ("A4 gets first-hand verification pain.", "A field coordinator describes a real cost from unverified / mis-verified information today."),
    ("Convergence, not a single source.", "Three independent sources agreeing beats any one report or interview — that is the bar for \u201cvalidated.\u201d"),
], size=14.5, gap=11)
footer(s)

# ============================================================ SLIDE 15 — shortest path to pilot
s = slide(); header(s, "Path", "14 · Shortest path to pilot", "14 / 15")
steps = [
    ("1", "Land Row-1", "Mones / Abdelrahman reply → route via evidence pipeline → A6-FIELD-LOG Row 1.", TEAL),
    ("2", "Close the two Low A6b sub-claims", "One first-hand informal-COLLAPSE (ERRs/Omar) + one informal-SUCCESS (White Helmets local team).", BLUE),
    ("3", "First A7 + A6-trust signal", "One adoption signal + one named-verifier signal from a real org/coordinator.", AMBER),
    ("4", "Phase-0 baseline + Phase-1 real data", "Lock prototype (tests) + HDX real seed — already scoped in PLAN.md.", GREEN),
    ("5", "Phase-3 auth + gated verify", "Only trusted actors verify — closes the security gap before any pilot.", NAVY),
]
y = Inches(1.5)
for n, t, d, col in steps:
    rect(s, Inches(0.55), y, Inches(0.7), Inches(0.92), col)
    txt(s, Inches(0.55), y+Inches(0.22), Inches(0.7), Inches(0.5),
        [[(n, 22, True, WHITE)]], align=PP_ALIGN.CENTER)
    rect(s, Inches(1.33), y, Inches(11.4), Inches(0.92), SAND)
    txt(s, Inches(1.5), y+Inches(0.08), Inches(11.1), Inches(0.4),
        [[(t, 14.5, True, NAVY)]])
    txt(s, Inches(1.5), y+Inches(0.47), Inches(11.1), Inches(0.4),
        [[(d, 12, False, INK)]])
    y += Inches(1.04)
footer(s)

# ============================================================ SLIDE 16 — what happens next
s = slide(); header(s, "Next", "15 · What happens next", "15 / 15")
bullet(s, Inches(0.55), Inches(1.5), Inches(12.2), Inches(3.2), [
    ("Hermes rests until real-world evidence arrives.", "Either Mones replies, or a new coordinator conversation starts. No new research, no outreach expansion, no features while pending replies are unresolved."),
    ("When a reply lands — existing pipeline only.", "EVIDENCE-INTAKE-TEMPLATE → A6-FIELD-LOG Row 1 → ASSUMPTION-TRACKER → journal → commit. A reply alone is never evidence."),
    ("Candidate list is deliberately frozen.", "Expand Tier-2/3 sourcing only after current Tier-1 replies resolve. No single-dependency on one person."),
], size=14.5, gap=11)
rect(s, Inches(0.55), Inches(5.0), Inches(12.2), Inches(1.35), NAVY)
txt(s, Inches(0.8), Inches(5.15), Inches(11.7), Inches(1.1),
    [[("Bottom line:", 15, True, TEAL),
      ("  The capability question (can structure survive?) is largely answered — conditionally. The open risk is adoption. The next useful action is processing one real coordinator story, not building more.", 14, False, WHITE)]],
    space_after=4)
footer(s)

# ============================================================ APPENDIX SLIDES
def appendix(title, body_runs):
    s = slide()
    rect(s, 0, 0, SW, SH, NAVY)
    rect(s, 0, Inches(0.0), SW, Inches(0.12), TEAL)
    txt(s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.6),
        [[("APPENDIX", 13, True, TEAL)]])
    txt(s, Inches(0.7), Inches(0.95), Inches(12), Inches(0.8),
        [[(title, 24, True, WHITE)]])
    txt(s, Inches(0.7), Inches(2.0), Inches(12.0), Inches(4.8),
        body_runs, space_after=10, line_spacing=1.05)
    return s

appendix("60-second pitch",
    [[("Nidaa is an offline-first, Arabic-RTL aid-coordination board for Syria and Gaza — built so it still works when the internet and power don't. ", 16, False, SAND),
      ("Connectivity is the failure mode in both places, and the tools people rely on (WhatsApp, Telegram) go dark exactly during the outages when aid info matters most. ", 16, False, SAND),
      ("Nidaa saves posts on the device, syncs when a connection appears, and is verifiable by trusted local organizations. ", 16, False, SAND),
      ("We've moved from \u201cwould this work?\u201d to a sharper question — and the evidence says structure survives only when it's near-zero-burden and owned by someone on the ground. That's the product we're building.", 16, True, WHITE)]])

appendix("5-minute explanation",
    [[("Start with the constraint: in Syria ~64% are offline; in Gaza, internet and electricity are scarce. Aid information fragments precisely when it's needed. ", 15, False, SAND),
      ("Existing tools assume a connection that isn't there. ", 15, False, SAND)],
     [("What Nidaa is: an offline-first board — the phone is the source of truth, posts queue locally and sync on reconnect, and trusted orgs mark entries verified. Arabic-default, full RTL, Gaza-first. ", 15, False, SAND)],
     [("What we learned: we started assuming a standalone coordinator-led board. The evidence (ACAPS on ERRs, White Helmets reports, GTS, CDAC) shows a cleaner truth — structure is maintained only where there's an owner and near-zero burden. ", 15, False, SAND),
      ("That's Reality B: communication survives, structure doesn't — unless you embed in the workflow and become the facilitating function. ", 15, True, WHITE)],
     [("Where we are: A6b is resolved-as-conditional; A2/A3/A5 strengthened; A7 (adoption) is the real remaining risk. ", 15, False, SAND),
      ("Next: one real coordinator story (Row-1), then close the adoption and trust questions before any pilot.", 15, True, WHITE)]])

appendix("Judge / competition explanation",
    [[("Category: offline-first humanitarian coordination tooling. ", 15, True, WHITE),
      ("We scanned the landscape (Jul 2026): crisis-mesh messengers exist, preparedness tools exist, but a needs-board that is genuinely offline-first AND Arabic-RTL is absent — globally, not just regionally. ", 15, False, SAND)],
     [("Our edge is the combination, not any single part: offline-first storage + verifiable entries + Arabic-RTL + region-tagging + (planned) device-to-device mesh sync and decentralized trust. ", 15, False, SAND)],
     [("The defensible claim: first-in-combination. Execution and field-validation are the real risks — which is exactly why the current work is evidence, not features. ", 15, True, WHITE)],
     [("Traction to date: working prototype (build green, API + Arabic round-trip verified), real-seed-data importer scoped, and a live evidence campaign that has already changed the product shape.", 15, False, SAND)]])

appendix("Humanitarian practitioner explanation",
    [[("If you coordinate aid on the ground: Nidaa is a board your team can post to and read even when the network is down, that trusted orgs verify, and that doesn't add admin burden to your day. ", 15, True, WHITE)],
     [("We're not asking you to adopt a new system — we're asking: when things got hard, who actually kept the information, and what happened to it? ", 15, False, SAND),
      ("That's the six-question story we need. ", 15, False, SAND)],
     [("Early signal says: where there's a named owner and the capture is trivial, structure holds (White Helmets log field + monthly reports under bombardment). Where it's a loose chat with no owner, documentation is the gap everyone names but nobody fills. ", 15, False, SAND)],
     [("So the design bet is: sit inside the channel you already use, make posting near-invisible, and let a facilitating function own verification — rather than hoping a separate board gets maintained. ", 15, True, WHITE)],
     [("If you've lived the collapse of a spreadsheet under surge, or the survival of one under fire — that's the evidence that shapes this.", 15, False, SAND, True)]])

prs.save("Nidaa-Founder-Briefing-v1.pptx")
print("OK -> Nidaa-Founder-Briefing-v1.pptx  (%d slides)" % len(prs.slides._sldIdLst))
